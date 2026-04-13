#!/usr/bin/env python3
"""Generate pipelight presentation PPT in Chinese."""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG = RGBColor(0x1B, 0x1B, 0x2F)
BLUE = RGBColor(0x00, 0xD2, 0xFF)
GREEN = RGBColor(0x00, 0xE6, 0x96)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)
RED = RGBColor(0xFF, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DIM = RGBColor(0x77, 0x77, 0x88)

FONT_CN = "PingFang SC"
FONT_CODE = "Menlo"


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def tb(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False,
       align=PP_ALIGN.LEFT, font=FONT_CN):
    """Add a textbox with a single run."""
    box = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf


def add_p(tf, text, sz=18, color=WHITE, bold=False, font=FONT_CN,
          space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.space_before = Pt(space_before)
    return p


def rect(slide, l, t, w, h, fill, text="", sz=14, color=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Cm(l), Cm(t), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if text:
        tf = s.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = True
        p.font.name = FONT_CN
        p.alignment = PP_ALIGN.CENTER
    return s


# ================================================================
prs = Presentation()
prs.slide_width = Cm(33.867)   # 16:9
prs.slide_height = Cm(19.05)

# ================================================================
# 1. Title
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 2, 4, 30, 3, "PIPELIGHT", 56, BLUE, True, PP_ALIGN.CENTER, FONT_CODE)
tb(s, 2, 7.5, 30, 1.5, "AI 原生的轻量级 CLI CI/CD 工具", 28, WHITE, False, PP_ALIGN.CENTER)
tb(s, 2, 9.5, 30, 1, "本地优先  |  Docker 隔离  |  LLM Harness 工程",
   18, GREEN, False, PP_ALIGN.CENTER)
tb(s, 2, 14, 30, 1, "Rust + Tokio + Bollard + Claude Code",
   14, DIM, False, PP_ALIGN.CENTER, FONT_CODE)

# ================================================================
# 2. Why Pipelight
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 1.5, 1, 30, 1.5, "为什么需要 Pipelight?", 36, BLUE, True)

# Left: traditional
tb(s, 1.5, 3.2, 14, 1, "传统 CI/CD", 22, ORANGE, True)
tf = tb(s, 1.5, 4.5, 14, 8,
        "构建失败 -> 输出错误日志", 16, GRAY)
add_p(tf, "-> 人工阅读日志", 16, GRAY)
add_p(tf, "-> 人工修复代码", 16, GRAY)
add_p(tf, "-> 手动重新触发流水线", 16, GRAY)
add_p(tf, "", 10, GRAY)
add_p(tf, "笨重、依赖云端、对执行内容毫无理解", 16, DIM)

# Right: pipelight
tb(s, 18, 3.2, 14, 1, "Pipelight", 22, GREEN, True)
tf = tb(s, 18, 4.5, 14, 8,
        "构建失败 -> Claude 分析错误", 16, WHITE)
add_p(tf, "-> Claude 自动修复代码", 16, WHITE)
add_p(tf, "-> 自动重新构建 -> 通过", 16, WHITE)
add_p(tf, "", 10, WHITE)
add_p(tf, "本地优先、CLI 原生、AI 原生", 16, GREEN, True)
add_p(tf, "LLM 是受控的执行单元，不是事后补丁", 16, GREEN)

# ================================================================
# 3. Architecture
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 1.5, 1, 30, 1.5, "六层流水线架构", 36, BLUE, True)

layers = [
    ("CLI (clap)", "解析命令: run / init / retry / status / clean", BLUE),
    ("Detector 检测器", "策略模式: 检测项目类型、提取版本信息", GREEN),
    ("PipelineBuilder 构建器", "StepDef trait: 生成步骤、依赖、失败策略", ORANGE),
    ("Parser 解析器", "YAML -> Pipeline 结构体, 校验 DAG 约束", RGBColor(0xDE, 0x5C, 0x0B)),
    ("Scheduler 调度器", "petgraph DAG -> 拓扑排序 -> 批次调度", BLUE),
    ("Executor 执行器", "bollard Docker API: 创建容器、挂载目录、收集日志", GREEN),
]

y = 3.5
for name, desc, color in layers:
    rect(s, 1.5, y, 8, 1.3, color, name, 14)
    tb(s, 10.5, y + 0.15, 22, 1.2, desc, 15, GRAY, font=FONT_CN)
    y += 1.6

# Data flow
tb(s, 1.5, y + 0.5, 30, 1,
   "init:  目录 -> Detector -> ProjectInfo -> Builder -> pipeline.yml",
   13, DIM, font=FONT_CODE)
tb(s, 1.5, y + 1.2, 30, 1,
   "run:   pipeline.yml -> Parser -> Scheduler(DAG) -> Executor(Docker) -> Output",
   13, DIM, font=FONT_CODE)

# ================================================================
# 4. Usage
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 1.5, 1, 30, 1.5, "核心命令", 36, BLUE, True)

cmds = [
    ("pipelight init", "检测项目，生成 pipeline.yml"),
    ("pipelight run", "在 Docker 容器中执行流水线"),
    ("pipelight retry --step <name>", "重试失败的步骤（同一 run-id）"),
    ("pipelight status --run-id <id>", "查看运行状态（tty/plain/json）"),
    ("pipelight list", "列出所有步骤和依赖关系"),
    ("pipelight validate", "校验 pipeline.yml 语法和 DAG"),
    ("pipelight clean", "清除项目产物"),
]

y = 3.5
for cmd, desc in cmds:
    tb(s, 2, y, 15, 1, f"$ {cmd}", 15, GREEN, True, font=FONT_CODE)
    tb(s, 18, y, 14, 1, desc, 15, GRAY)
    y += 1.5

tb(s, 2, y + 0.3, 30, 1,
   "输出模式:  --output tty（终端彩色）  |  plain（纯文本）  |  json（给 LLM）",
   14, ORANGE)

# ================================================================
# 5. pipeline.yml
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 1.5, 1, 30, 1.5, "pipeline.yml: 唯一真相来源", 36, BLUE, True)

yaml_text = (
    "name: gradle-java-ci\n"
    "image: gradle:8-jdk17\n"
    "steps:\n"
    "  - name: build\n"
    "    commands:\n"
    "      - gradle assemble\n"
    "    on_failure:\n"
    "      callback_command: abort\n"
    "\n"
    "  - name: pmd\n"
    "    depends_on: [build]\n"
    "    on_failure:\n"
    "      callback_command: auto_fix\n"
    "      max_retries: 2\n"
    "      exceptions:\n"
    "        ruleset_not_found:\n"
    "          command: auto_gen_pmd_ruleset\n"
    "          max_retries: 2"
)
tb(s, 1.5, 3, 16, 14, yaml_text, 12, GREEN, font=FONT_CODE)

# Right: key fields
fields = [
    ("commands", "Docker 容器内执行的 shell 命令"),
    ("on_failure", "步骤失败后的处理策略"),
    ("callback_command", "发给 LLM 的指令（枚举值）"),
    ("max_retries", "重试上限（熔断器）"),
    ("exceptions", "stderr 模式 -> 特定处理器"),
    ("context_paths", "LLM 被允许读取的文件列表"),
]
y = 3.2
for key, val in fields:
    tb(s, 19, y, 13, 0.7, key, 15, ORANGE, True, font=FONT_CODE)
    tb(s, 19, y + 0.8, 13, 0.7, val, 14, GRAY)
    y += 2

# ================================================================
# 6. Harness Engineering
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 1.5, 1, 30, 1.5, "Harness Engineering 驾驭工程", 36, BLUE, True)

tb(s, 1.5, 3, 30, 1,
   "Harness = 包裹在非确定性 LLM 外层的确定性控制平面", 20, GREEN, True)

tb(s, 1.5, 4.3, 30, 1.2,
   '"让 LLM 更聪明"是次要的，"让 LLM 嵌入一个永不失控的 harness"才是工程化重点。',
   16, ORANGE, True)

principles = [
    ("确定性控制平面", "pipelight（Rust）状态机 + DAG 调度决定流程走向，LLM 零自主权"),
    ("有界自治", "每个 CallbackCommand 带 max_retries，配额耗尽即 fail-close，无死循环"),
    ("显式契约", "pipelight <-> LLM 接口是严格 JSON + CallbackCommand 枚举，无自由文本"),
    ("职责隔离", "容器跑命令 / Harness 做决策 / LLM 做修复——三层互不越权"),
    ("最小权限", "LLM 只能读 context_paths 列出的文件，只能通过 pipelight retry 触发重试"),
    ("失败分类", "CallbackCommand 枚举把失败分为 retry / skip / abort 三族，非「看着办」"),
]

y = 6
for title, desc in principles:
    tb(s, 2, y, 6, 1, title, 16, BLUE, True)
    tb(s, 8.5, y, 24, 1, desc, 14, GRAY)
    y += 1.7

# ================================================================
# 7. Three-Layer Topology
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 1.5, 1, 30, 1.5, "三层 Harness 拓扑", 36, BLUE, True)

# Docker
rect(s, 1.5, 3.5, 9.5, 4, RGBColor(0x1E, 0x3A, 0x5F))
tb(s, 2, 3.7, 8.5, 1, "Docker 容器", 20, BLUE, True, PP_ALIGN.CENTER)
tf = tb(s, 2, 5, 8.5, 2, "执行 shell 命令", 14, GRAY, align=PP_ALIGN.CENTER)
add_p(tf, "产出 stdout / stderr / exit_code", 14, GRAY)
add_p(tf, "不做任何决策，只跑完返回退出码", 14, GRAY)

# Pipelight
rect(s, 12.3, 3.5, 9.5, 4, RGBColor(0x3A, 0x2A, 0x0A))
tb(s, 12.8, 3.7, 8.5, 1, "Pipelight (Rust)", 20, ORANGE, True, PP_ALIGN.CENTER)
tf = tb(s, 12.8, 5, 8.5, 2, "match_exception(stderr)", 14, GRAY, align=PP_ALIGN.CENTER, font=FONT_CODE)
add_p(tf, "查 pipeline.yml 的 on_failure", 14, GRAY, font=FONT_CN)
add_p(tf, "决定: retry / skip / abort", 14, GRAY, font=FONT_CN)
add_p(tf, "输出结构化 JSON", 14, GRAY, font=FONT_CN)

# LLM
rect(s, 23.1, 3.5, 9.5, 4, RGBColor(0x0A, 0x3A, 0x2A))
tb(s, 23.6, 3.7, 8.5, 1, "LLM (Claude)", 20, GREEN, True, PP_ALIGN.CENTER)
tf = tb(s, 23.6, 5, 8.5, 2, "读 on_failure.command", 14, GRAY, align=PP_ALIGN.CENTER)
add_p(tf, "查 skill 回调命令处理表", 14, GRAY)
add_p(tf, "执行: 修代码 / 生成配置", 14, GRAY)
add_p(tf, "调用 pipelight retry", 14, GRAY)

# Arrows
tb(s, 1.5, 8, 30, 1,
   "exit_code + stderr  ------->  CallbackCommand (JSON)  ------->  具体动作 + retry",
   14, DIM, False, PP_ALIGN.CENTER, FONT_CODE)

# Boundary rules
tb(s, 1.5, 9.5, 30, 1, "边界划分:", 20, ORANGE, True)
rules = [
    "Docker 容器回答:「命令跑完了没？退出码是多少？」 (退出码规则写在 pipeline.yml)",
    "Pipelight 回答:「失败了该干什么？」(决策表写在 pipeline.yml，Rust 只做 stderr -> exception_key 翻译)",
    "LLM 回答:「按指令具体怎么干？」 (动作定义在 skill 的回调命令处理表)",
]
y = 11
for r in rules:
    tb(s, 2, y, 30, 1, r, 13, GRAY)
    y += 1.3

# ================================================================
# 8. CallbackCommand
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 1.5, 1, 30, 1.5, "CallbackCommand: Harness 的指令集", 36, BLUE, True)
tb(s, 1.5, 2.7, 30, 1,
   "可穷举 | 可审计 | 可扩展 —— 永远没有自由文本指令", 16, GREEN)

# Table header
hy = 4.2
tb(s, 1.5, hy, 5.5, 0.8, "指令", 14, ORANGE, True)
tb(s, 7, hy, 3, 0.8, "动作", 14, ORANGE, True)
tb(s, 10, hy, 3.5, 0.8, "流水线状态", 14, ORANGE, True)
tb(s, 13.5, hy, 10, 0.8, "LLM 执行什么", 14, ORANGE, True)
tb(s, 24, hy, 8, 0.8, "典型场景", 14, ORANGE, True)

cmds = [
    ("AutoFix", "Retry", "retryable", "读 stderr + 修复源码", "编译/测试失败", GREEN),
    ("AutoGenPmdRuleset", "Retry", "retryable", "搜索/生成 PMD 规则集", "PMD 规则缺失", GREEN),
    ("Ping", "Retry", "retryable", "打印 pong，循环 10 轮", "连通性测试", GREEN),
    ("GitFail", "Skip", "继续", "无操作（自动跳过）", "网络/认证失败", BLUE),
    ("FailAndSkip", "Skip", "继续", "无操作（自动跳过）", "前置条件缺失", BLUE),
    ("RuntimeError", "Abort", "failed", "报告错误，终止", "运行时错误", RED),
    ("Abort", "Abort", "failed", "报告错误，终止", "严重代码问题", RED),
]

y = 5.2
for name, action, status, llm_do, scenario, color in cmds:
    tb(s, 1.5, y, 5.5, 0.8, name, 12, color, True, font=FONT_CODE)
    tb(s, 7, y, 3, 0.8, action, 12, GRAY)
    tb(s, 10, y, 3.5, 0.8, status, 12, GRAY)
    tb(s, 13.5, y, 10, 0.8, llm_do, 12, GRAY)
    tb(s, 24, y, 8, 0.8, scenario, 12, DIM)
    y += 1.2

# Three families
tb(s, 1.5, y + 0.5, 30, 1, "三族失败响应策略:", 18, ORANGE, True)
fams = [
    ("Retry 族", "有界自治修复。max_retries 是熔断器。", GREEN),
    ("Skip 族", "降级继续。Harness 自主决策，LLM 完全不介入。", BLUE),
    ("Abort 族", "快速失败。保护下游不受污染数据影响。", RED),
]
y2 = y + 1.8
for fn, fd, fc in fams:
    tb(s, 2, y2, 4, 0.8, fn, 15, fc, True)
    tb(s, 6.5, y2, 26, 0.8, fd, 14, GRAY)
    y2 += 1.1

# ================================================================
# 9. Real Example
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 1.5, 1, 30, 1.5, "实战: PMD 规则集自动生成", 36, BLUE, True)

steps_ex = [
    ("1", "容器执行 PMD 步骤",
     "未找到 pmd-ruleset.xml -> 输出 PIPELIGHT_CALLBACK:auto_gen_pmd_ruleset -> exit 1",
     BLUE),
    ("2", "Pipelight 输出 JSON",
     "match_exception -> ruleset_not_found -> command: auto_gen_pmd_ruleset -> status: retryable",
     ORANGE),
    ("3", "LLM 执行回调",
     "第一轮: 搜索项目中已有 PMD 配置 | 第二轮: 搜索编码规范文档 -> 生成 pmd-ruleset.xml",
     GREEN),
    ("4", "LLM 触发重试",
     "pipelight retry --run-id run-001 --step pmd（同一 run-id 关联同一执行上下文）",
     GREEN),
    ("5", "容器重新执行",
     "pmd-ruleset.xml 存在 -> 执行 PMD 扫描 -> PMD Total: 1604 violations -> exit 0",
     BLUE),
    ("6", "最终 JSON 输出",
     "steps[pmd].status = success，流水线继续执行后续步骤",
     ORANGE),
]

y = 3.2
for num, title, desc, color in steps_ex:
    rect(s, 1.5, y, 1.3, 1.3, color, num, 18)
    tb(s, 3.3, y, 7, 1.3, title, 15, color, True)
    tb(s, 10.5, y, 22, 1.3, desc, 13, GRAY)
    y += 1.8

# Key points
tb(s, 1.5, y + 0.3, 30, 1, "Harness 工程要点:", 16, ORANGE, True)
obs = [
    "pipelight 不告诉 LLM「去搜 docs/」——只发抽象指令 auto_gen_pmd_ruleset，具体怎么搜由 skill 定义",
    "retries_remaining=2 是硬上限，两次失败后 skill 强制 --skip pmd 新起 run，无死循环",
    "run_id 贯穿全程，初跑 + retry 绑定为同一执行上下文，幂等、可重放、可审计",
]
y2 = y + 1.3
for o in obs:
    tb(s, 2, y2, 30, 0.9, o, 12, GRAY)
    y2 += 1

# ================================================================
# 10. Control Flow
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 1.5, 1, 30, 1.5, "完整控制链路", 36, BLUE, True)

chain = [
    ("[Docker 容器]", "执行 pipeline.yml 的 commands，产出 stdout/stderr/exit_code", BLUE),
    ("[Executor 层]", "捕获退出码和输出流，落盘日志到 pipelight-misc/", DIM),
    ("[StepDef::match_exception()]", "用 stderr 特征串匹配 exception_key", ORANGE),
    ("[查 pipeline.yml]", "exception_key -> { command, max_retries, context_paths }", ORANGE),
    ("[CallbackCommandRegistry]", "CallbackCommand -> Action (Retry / Skip / Abort)", GREEN),
    ("[JSON 输出]", "结构化输出: status + on_failure.command + retries_remaining", GREEN),
    ("[LLM + Skill]", "按回调命令处理表执行动作，调用 pipelight retry 回到第一步", BLUE),
]

y = 3.2
for i, (label, desc, color) in enumerate(chain):
    tb(s, 2, y, 10, 1, label, 14, color, True, font=FONT_CODE)
    tb(s, 13, y, 20, 1, desc, 14, GRAY)
    if i < len(chain) - 1:
        tb(s, 6, y + 0.9, 2, 0.6, "|", 14, DIM, False, PP_ALIGN.CENTER, FONT_CODE)
        tb(s, 6.5, y + 0.9, 2, 0.6, "v", 12, DIM, False, PP_ALIGN.LEFT, FONT_CODE)
    y += 1.6

# ================================================================
# 11. Design Philosophy
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 1.5, 1, 30, 1.5, "设计哲学", 36, BLUE, True)

phils = [
    ("配置即契约", "pipeline.yml 是完整的运行时规格。任何工程师都能在不读 Rust 源码的\n前提下审计、修改、回滚 harness 的决策。", GREEN),
    ("声明式指令", "Harness 发出「要什么」(auto_gen_pmd_ruleset)，而非「怎么做」(搜 docs/)。\n「怎么做」在 skill 中定义，与 pipelight 完全解耦。", BLUE),
    ("模板与运行时分离", "Rust 代码 = 烘焙模板（通过 pipelight init 生成 pipeline.yml）\npipeline.yml = 运行时契约（pipelight run 只读它）", ORANGE),
    ("人在回路", "pipeline.yml 对用户完全透明可编辑。\n任何 harness 决策都能在 YAML 中核对和覆盖。", RGBColor(0xDE, 0x5C, 0x0B)),
    ("每层可替换", "换 Docker 为 Podman、换调度器、换 Claude 为其他 LLM——\n层间的 JSON 契约保持稳定。", GREEN),
]

y = 3.3
for title, desc, color in phils:
    tb(s, 2, y, 7, 1, title, 18, color, True)
    tb(s, 2, y + 1.2, 30, 2, desc, 14, GRAY)
    y += 2.8

# ================================================================
# 12. Tech Stack
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 1.5, 1, 30, 1.5, "技术栈", 36, BLUE, True)

stack = [
    ("Rust", "核心语言 — 安全、高性能、零成本抽象"),
    ("clap", "CLI 参数解析 — run / init / retry / status / list / validate / clean"),
    ("tokio", "异步运行时 — 并行步骤执行、Docker 日志流"),
    ("bollard", "Docker API 客户端 — 不通过 shell 调用 docker CLI"),
    ("petgraph", "DAG 构建和拓扑排序，用于步骤调度"),
    ("serde + serde_yaml", "流水线 YAML 序列化 / 反序列化"),
    ("anyhow + thiserror", "错误处理: anyhow（应用层）+ thiserror（库层）"),
    ("tracing", "结构化日志（不用 println! 或 log crate）"),
    ("indicatif + console", "终端进度条和彩色输出"),
]

y = 3.3
for name, desc in stack:
    tb(s, 2, y, 6, 0.9, name, 16, ORANGE, True, font=FONT_CODE)
    tb(s, 9, y, 24, 0.9, desc, 14, GRAY)
    y += 1.5

# ================================================================
# 13. Summary
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tb(s, 2, 3, 30, 2, "PIPELIGHT", 48, BLUE, True, PP_ALIGN.CENTER, FONT_CODE)

items = [
    "本地优先 CLI CI/CD — 无需服务器，无需云端依赖",
    "Docker 隔离 — 每个步骤在容器内执行",
    "DAG 调度 — 无依赖步骤并行执行",
    "AI 原生 — LLM 作为受控执行单元，Harness Engineering 驾驭",
    "故障安全 — 有界自治、熔断器、fail-close 默认策略",
    "可审计 — 结构化 JSON、持久化运行状态、可重放",
]

y = 6
for item in items:
    tb(s, 4, y, 26, 1, item, 17, GRAY, False, PP_ALIGN.LEFT)
    y += 1.5

tb(s, 2, 16, 30, 1,
   "github.com/amwtke/deeparch_cicd", 14, DIM, False, PP_ALIGN.CENTER, FONT_CODE)

# ================================================================
output = "/Users/xiaojin/workshop/deeparch_cicd/docs/pipelight-presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
